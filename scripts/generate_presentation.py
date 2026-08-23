"""Generate doc/songstem_overview.pptx — a 5-slide, diagram-first overview of Songstem.

Run after changing the deck content:  python scripts/generate_presentation.py
Requires python-pptx (in the dev extras: pip install -e ".[dev]").
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

TARGET = Path(__file__).resolve().parents[1] / "doc" / "songstem_overview.pptx"

# ---------------------------------------------------------------- palette
INK = RGBColor(0x1F, 0x29, 0x37)      # near-black text
MUTED = RGBColor(0x6B, 0x72, 0x80)    # gray text
BLUE = RGBColor(0x25, 0x63, 0xEB)     # primary accent
BLUE_FILL = RGBColor(0xEF, 0xF6, 0xFF)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_FILL = RGBColor(0xF0, 0xFD, 0xFA)
AMBER = RGBColor(0xB4, 0x53, 0x09)
AMBER_FILL = RGBColor(0xFF, 0xFB, 0xEB)
GRAY_LINE = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_FILL = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _set_text(shape, lines, size=13, color=INK, align=PP_ALIGN.CENTER, bold_first=True):
    """Fill a shape's text frame: `lines` is a list of (text, overrides) or plain strings."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(6)
    tf.margin_top = tf.margin_bottom = Pt(4)
    for i, line in enumerate(lines):
        text_, over = (line, {}) if isinstance(line, str) else line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text_
        f = run.font
        f.name = FONT
        f.size = Pt(over.get("size", size))
        f.color.rgb = over.get("color", color)
        f.bold = over.get("bold", bold_first and i == 0)


def box(slide, x, y, w, h, lines, fill=BLUE_FILL, line=BLUE, size=13, color=INK,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.CENTER):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = 0.12
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line
    sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    _set_text(sp, lines, size=size, color=color, align=align)
    return sp


def arrow(slide, x, y, w=0.4, h=0.28, direction="right"):
    shape = {"right": MSO_SHAPE.RIGHT_ARROW, "down": MSO_SHAPE.DOWN_ARROW,
             "left": MSO_SHAPE.LEFT_ARROW}[direction]
    if direction == "down":
        w, h = h, w
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = GRAY_LINE
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, lines, size=13, color=INK, align=PP_ALIGN.LEFT, bold_first=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set_text(tb, lines, size=size, color=color, align=align, bold_first=bold_first)
    return tb


def title(slide, txt, subtitle=None):
    text(slide, 0.55, 0.28, 12.2, 0.7, [(txt, {"size": 30, "bold": True})], align=PP_ALIGN.LEFT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.98),
                                 Inches(1.6), Pt(3.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False
    if subtitle:
        text(slide, 0.55, 1.08, 12.2, 0.4, [(subtitle, {"size": 14, "color": MUTED})])


# ================================================================ slide 1 — title
s = prs.slides.add_slide(BLANK)
text(s, 0.9, 1.15, 11.5, 1.0, [("Songstem", {"size": 54, "bold": True, "color": BLUE})])
text(s, 0.9, 2.15, 11.5, 0.5,
     [("Instrument removal for Apple Music on Windows — batch stem separation "
       "for music practice", {"size": 19, "color": INK})])

# concept diagram: full mix -> solo + muted
box(s, 1.6, 3.7, 2.9, 1.15, ["Full mix", ("any song in an iTunes playlist",
                                          {"size": 11, "color": MUTED})],
    fill=GRAY_FILL, line=GRAY_LINE, size=16)
arrow(s, 4.75, 3.85, 0.55)
arrow(s, 4.75, 4.75, 0.55)
box(s, 5.55, 3.45, 3.4, 0.95, ["Solo", ("the chosen stem, isolated — e.g. bass only",
                                        {"size": 11, "color": MUTED})], size=16)
box(s, 5.55, 4.6, 3.4, 0.95, ["Muted", ("the mix with that stem removed",
                                        {"size": 11, "color": MUTED})],
    fill=TEAL_FILL, line=TEAL, size=16)
text(s, 9.3, 3.9, 3.2, 1.3,
     [("Play along with the band —", {"size": 13, "color": MUTED}),
      ("or study your part in isolation.", {"size": 13, "color": MUTED})])
text(s, 0.9, 6.6, 11.5, 0.4,
     [("Windows · Python · optimized for practice and cover-song rehearsal — "
       "quality over latency (batch, not live)", {"size": 12, "color": MUTED})])

# ================================================================ slide 2 — usage model
s = prs.slides.add_slide(BLANK)
title(s, "Usage model", "One button after setup — the batch runs in the background")

steps = [
    ("1 · Pick a playlist", "read live from iTunes /\nApple Music (COM)"),
    ("2 · Check songs", "unusable (DRM, cloud-only)\ntracks are greyed out"),
    ("3 · Choose the stem", "bass, guitar, vocals… plus\nper-stem muted-mix levels"),
    ("4 · Run", "separation runs in a child\nprocess; progress + Stop"),
]
x = 0.6
for head, detail in steps:
    lines = [head] + [(ln, {"size": 11, "color": MUTED}) for ln in detail.split("\n")]
    box(s, x, 2.0, 2.55, 1.35, lines, size=15)
    x += 2.55
    if head[0] != "4":
        arrow(s, x + 0.05, 2.55)
        x += 0.5

arrow(s, 11.0, 3.55, direction="down", w=0.6)

# outputs row (header anchored to the top so the file boxes below don't cover it)
outputs = box(s, 8.05, 4.35, 4.7, 1.85,
              [("Per song, in the output folder", {"size": 13, "bold": True, "color": TEAL})],
              fill=TEAL_FILL, line=TEAL)
outputs.text_frame.vertical_anchor = MSO_ANCHOR.TOP
oy = 4.85
for name in ["Artist – Title [bass solo].wav",
             "Artist – Title [bass muted].wav",
             "Artist – Title [bass cheatsheet].md"]:
    box(s, 8.3, oy, 4.2, 0.38, [(name, {"size": 11, "bold": False})],
        fill=WHITE, line=TEAL, size=11)
    oy += 0.44
text(s, 8.55, 6.25, 4.2, 0.3, [("cheat sheet: key, tempo, note patterns, synced lyrics",
                                {"size": 10.5, "color": MUTED})])

arrow(s, 7.35, 5.1, 0.55, direction="left")  # outputs feed the practice step
box(s, 4.4, 4.6, 2.85, 1.35,
    ["5 · Practice", ("audition solo / muted in the", {"size": 11, "color": MUTED}),
     ("built-in player; selections are", {"size": 11, "color": MUTED}),
     ("remembered between launches", {"size": 11, "color": MUTED})], size=15)
text(s, 0.6, 4.75, 3.6, 1.2,
     [("Optional per-song pitch shift", {"size": 12, "bold": True}),
      ("±12 half-steps → transposed WAV copies,", {"size": 11, "color": MUTED}),
      ("for comfortable keys / alternate tunings", {"size": 11, "color": MUTED})],
     bold_first=True)

# ================================================================ slide 3 — sources & DRM
s = prs.slides.add_slide(BLANK)
title(s, "Audio sources and the DRM boundary",
      "Separation needs decodable audio — Apple's FairPlay DRM cannot be decrypted")

box(s, 4.7, 1.75, 3.9, 0.85, ["iTunes / Apple Music track"],
    fill=GRAY_FILL, line=GRAY_LINE, size=15)
arrow(s, 5.35, 2.75, direction="down", w=0.45)
arrow(s, 7.55, 2.75, direction="down", w=0.45)

# left branch: DRM-free
box(s, 1.3, 3.35, 4.6, 1.0,
    ["DRM-free file", ("CD rips, .mp3, DRM-free .m4a purchases",
                       {"size": 11, "color": MUTED})], size=15)
arrow(s, 3.35, 4.5, direction="down", w=0.45)

# right branch: DRM -> loopback
box(s, 7.4, 3.35, 4.6, 1.0,
    [("DRM-protected .m4p", {"color": AMBER}),
     ("subscription downloads — undecodable, greyed out", {"size": 11, "color": MUTED})],
    fill=AMBER_FILL, line=AMBER, size=15)
arrow(s, 9.5, 4.5, direction="down", w=0.45)
box(s, 6.55, 5.05, 6.3, 1.15,
    ["Loopback re-record → DRM-free WAV",
     ("iTunes plays  →  VB-Audio virtual cable  →  capture (sounddevice)",
      {"size": 11, "color": MUTED}),
     ("resumable · atomic writes · silent captures stop the run · personal use only",
      {"size": 10.5, "color": MUTED})],
    fill=AMBER_FILL, line=AMBER, size=14)
arrow(s, 5.95, 5.5, 0.5, direction="left")  # recordings feed back into separation

box(s, 1.3, 5.05, 4.1, 1.15,
    ["Stem separation", ("the resolved source — original file or", {"size": 11, "color": MUTED}),
     ("re-recorded WAV — feeds the pipeline", {"size": 11, "color": MUTED})],
    fill=TEAL_FILL, line=TEAL, size=15)
text(s, 0.6, 6.65, 12.2, 0.4,
     [("Re-recording captures playback in real time (a playlist takes about its runtime); "
       "it records audio output — it does not remove DRM.", {"size": 11.5, "color": MUTED})])

# ================================================================ slide 4 — architecture
s = prs.slides.add_slide(BLANK)
title(s, "Architecture", "Layered pipeline — every heavy dependency sits behind a swappable seam")

# GUI bar
box(s, 0.6, 1.7, 12.15, 0.75,
    [("GUI — PySide6 (Qt 6): playlist picker · stem & level controls · progress/Stop · "
      "built-in player (Qt Multimedia) · UI state in SQLite (pref)", {"size": 12.5})],
    fill=BLUE_FILL, line=BLUE)
arrow(s, 6.45, 2.55, direction="down", w=0.4)

# pipeline flow
flow = [
    ("Playlist source", "iTunes COM (pywin32)\nor a folder of WAVs", GRAY_FILL, GRAY_LINE),
    ("Batch pipeline", "child process — GUI\nnever stalls; per-song\nerrors don't abort",
     BLUE_FILL, BLUE),
    ("Separation backend", "Demucs (PyTorch, CPU)\npluggable via registry —\ne.g. Open-Unmix",
     TEAL_FILL, TEAL),
    ("Mixer", "pure NumPy —\nsolo & muted sums,\nper-stem gains", TEAL_FILL, TEAL),
    ("Audio I/O", "soundfile + ffmpeg\natomic writes\n→ WAVs on disk", GRAY_FILL, GRAY_LINE),
]
x = 0.6
for head, detail, fill, ln in flow:
    lines = [head] + [(t, {"size": 10.5, "color": MUTED}) for t in detail.split("\n")]
    box(s, x, 3.05, 2.2, 1.6, lines, fill=fill, line=ln, size=13.5)
    x += 2.2
    if head != "Audio I/O":
        arrow(s, x + 0.03, 3.7, 0.3)
        x += 0.36

arrow(s, 9.9, 4.8, direction="down", w=0.4)
box(s, 7.7, 5.35, 5.05, 1.15,
    ["Analysis (best-effort, never fails a job)",
     ("librosa: key (Krumhansl–Schmuckler) · tempo · note patterns",
      {"size": 11, "color": MUTED}),
     ("synced lyrics (network) → Markdown cheat sheet", {"size": 11, "color": MUTED})],
    size=13)
text(s, 0.6, 5.5, 6.6, 1.1,
     [("Design rule", {"size": 12.5, "bold": True}),
      ("Heavy/platform imports (Demucs, Qt, COM, soundfile) stay local behind ABCs "
       "with in-memory fakes — the core is importable and tested headless.",
       {"size": 11.5, "color": MUTED})], bold_first=True)

# ================================================================ slide 5 — technologies
s = prs.slides.add_slide(BLANK)
title(s, "Underlying technologies", "All Python — Windows x86 desktop, Python 3.14")

cards = [
    ("Separation / ML", TEAL, TEAL_FILL,
     ["Demucs v4 — hybrid transformer source separation",
      "PyTorch, CPU inference (GPU-ready via config)",
      "Backend registry → swappable engines"]),
    ("Desktop UI", BLUE, BLUE_FILL,
     ["PySide6 (Qt 6) widgets",
      "Qt Multimedia output player",
      "QThread workers + child process for heavy work"]),
    ("Audio processing", TEAL, TEAL_FILL,
     ["NumPy mixing · soundfile (libsndfile) I/O",
      "ffmpeg for AAC/ALAC decode (auto-installed via winget)",
      "librosa — key/tempo/note analysis, pitch shifting"]),
    ("Windows integration", BLUE, BLUE_FILL,
     ["pywin32 — iTunes/Apple Music COM automation",
      "sounddevice + VB-Audio Virtual Cable loopback capture",
      "Runs on regular CPython (not free-threaded 3.14t)"]),
    ("Persistence & state", AMBER, AMBER_FILL,
     ["pref — SQLite-backed UI state",
      "Remembers playlist, checked songs, pitch, window",
      "Atomic file writes for crash-safe outputs"]),
    ("Quality & tooling", AMBER, AMBER_FILL,
     ["pytest suite over the dependency-light core",
      "Fakes for iTunes, separation, recording, analysis",
      "ruff linting · hatchling packaging"]),
]
cw, ch, gx, gy = 4.05, 2.15, 0.15, 0.25
for i, (head, ln, fill, items) in enumerate(cards):
    cx = 0.6 + (i % 3) * (cw + gx)
    cy = 1.85 + (i // 3) * (ch + gy)
    lines = [(head, {"size": 14, "bold": True, "color": ln})]
    lines += [(f"•  {t}", {"size": 11, "color": INK}) for t in items]
    b = box(s, cx, cy, cw, ch, lines, fill=fill, line=ln, align=PP_ALIGN.LEFT)
    b.text_frame.margin_left = Pt(10)
text(s, 0.6, 6.55, 12.2, 0.4,
     [("Every dependency is plain PyPI — no services, no accounts; everything runs and "
       "stays on the local machine.", {"size": 12, "color": MUTED})])

TARGET.parent.mkdir(parents=True, exist_ok=True)
prs.save(TARGET)
print(f"wrote {TARGET}")
